from __future__ import annotations
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pydantic import BaseModel
from utils.fs import SLIDES_OUT, ts

# -----------------------
# Models-----
# --------------------------
class Bullet(BaseModel):
    text: str
    timestamp: str | None = None  # kept for future, but not shown in slides

class Section(BaseModel):
    heading: str
    bullets: list[Bullet]

class Outline(BaseModel):
    title: str
    topics: list[str]
    sections: list[Section]

TITLE_LAYOUT = 0
TITLE_AND_CONTENT = 1

# -----------------------
# Helpers
# -----------------------
def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])

    # Nude background color (light peach tone)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(222, 237, 255)  # color

    # Title formatting
    title_shape = slide.shapes.title
    title_shape.text = title[:120]
    for p in title_shape.text_frame.paragraphs:
        p.font.size = Pt(44)
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True

    # Subtitle formatting
    if slide.placeholders and len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle[:200]
        for p in sub.text_frame.paragraphs:
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True

def _add_content_slide(prs: Presentation, heading: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])

    # Heading formatting
    title_shape = slide.shapes.title
    title_shape.text = heading[:100]
    for p in title_shape.text_frame.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Bullets formatting
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0, 51, 102)  # dark blue
        # ⚠️ No auto-bold for first bullet now

def outline_to_pptx(outline: Outline, filename_stem: Optional[str] = None) -> Path:
    prs = Presentation()

    # Title slide
    _add_title_slide(prs, outline.title, ", ".join(outline.topics)[:200])

    # Content slides
    for sec in outline.sections:
        bullets = [item.text for item in sec.bullets]  # timestamps removed
        _add_content_slide(prs, sec.heading, bullets)

    # Save
    stem = filename_stem or f"slides_{ts()}"
    out_path = SLIDES_OUT / f"{stem}.pptx"
    prs.save(out_path)
    return out_path
